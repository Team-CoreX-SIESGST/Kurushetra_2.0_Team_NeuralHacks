import os
import json
import csv
import xml.etree.ElementTree as ET
from typing import Dict, Any, List, Union
import mimetypes
from datetime import datetime

# Import libraries for different file types
try:
    import PyPDF2
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
    # Configure tesseract path if provided or discover common Windows install locations
    
    # Optional lightweight transformer OCR (TrOCR base model)
    from transformers import TrOCRProcessor, VisionEncoderDecoderModel
    import torch
    TRANSFORMER_OCR_AVAILABLE = True
except ImportError:
    TRANSFORMER_OCR_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

# Optional hosted OCR via Hugging Face Inference API
try:
    from web_hf_inference import HFInferenceClient
    HF_INFERENCE_AVAILABLE = True
except Exception:
    HF_INFERENCE_AVAILABLE = False

class FileProcessor:
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx if DOCX_AVAILABLE else None,
            '.doc': self._process_docx if DOCX_AVAILABLE else None,
            '.xlsx': self._process_excel if PANDAS_AVAILABLE else None,
            '.xls': self._process_excel if PANDAS_AVAILABLE else None,
            '.csv': self._process_csv,
            '.txt': self._process_text,
            '.json': self._process_json,
            '.xml': self._process_xml,
            '.html': self._process_html if BS4_AVAILABLE else None,
            '.htm': self._process_html if BS4_AVAILABLE else None,
            '.md': self._process_markdown if MARKDOWN_AVAILABLE else None,
            '.pptx': self._process_pptx if PPTX_AVAILABLE else None,
            '.ppt': self._process_pptx if PPTX_AVAILABLE else None,
            '.png': self._process_image_ocr if OCR_AVAILABLE else None,
            '.jpg': self._process_image_ocr if OCR_AVAILABLE else None,
            '.jpeg': self._process_image_ocr if OCR_AVAILABLE else None,
            '.tiff': self._process_image_ocr if OCR_AVAILABLE else None,
        }
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported"""
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.supported_formats and self.supported_formats[ext] is not None
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return [ext for ext, processor in self.supported_formats.items() if processor is not None]
    
    def process_file(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        """Main method to process any supported file"""
        ext = os.path.splitext(original_filename.lower())[1]
        
        if not self.is_supported_format(original_filename):
            raise ValueError(f"Unsupported file format: {ext}")
        
        processor = self.supported_formats[ext]
        
        try:
            result = processor(file_path)
            
            # Add metadata to the result
            result.update({
                "file_metadata": {
                    "original_filename": original_filename,
                    "file_extension": ext,
                    "file_size": os.path.getsize(file_path),
                    "processed_at": datetime.now().isoformat(),
                    "processing_method": processor.__name__
                }
            })
            
            return result
            
        except Exception as e:
            raise Exception(f"Error processing {ext} file: {str(e)}")
    
    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process PDF files"""
        text_content = []
        metadata = {}
        
        try:
            # Try with pdfplumber first (better text extraction)
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                metadata = {
                    "pages": len(pdf.pages),
                    "metadata": pdf.metadata or {}
                }
                
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append({
                            "page": i + 1,
                            "content": page_text.strip()
                        })
                        
                        # Extract tables if any
                        tables = page.extract_tables()
                        if tables:
                            text_content[-1]["tables"] = tables
        
        except Exception as e:
            # Fallback to PyPDF2
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata = {
                        "pages": len(pdf_reader.pages),
                        "metadata": pdf_reader.metadata or {}
                    }
                    
                    for i, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_content.append({
                                "page": i + 1,
                                "content": page_text.strip()
                            })
            except Exception as fallback_error:
                raise Exception(f"PDF processing failed: {str(e)}, Fallback error: {str(fallback_error)}")
        
        return {
            "content_type": "pdf",
            "text_content": text_content,
            "document_metadata": metadata,
            "total_pages": len(text_content)
        }
    
    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process Word documents"""
        doc = Document(file_path)
        
        paragraphs = []
        tables = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append({
                    "text": para.text,
                    "style": para.style.name if para.style else "Normal"
                })
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
        
        # Extract document properties
        core_props = doc.core_properties
        metadata = {
            "author": core_props.author,
            "title": core_props.title,
            "subject": core_props.subject,
            "created": str(core_props.created) if core_props.created else None,
            "modified": str(core_props.modified) if core_props.modified else None
        }
        
        return {
            "content_type": "docx",
            "paragraphs": paragraphs,
            "tables": tables,
            "document_metadata": metadata,
            "total_paragraphs": len(paragraphs),
            "total_tables": len(tables)
        }
    
    def _process_excel(self, file_path: str) -> Dict[str, Any]:
        """Process Excel files"""
        excel_data = pd.ExcelFile(file_path)
        sheets_data = {}
        
        for sheet_name in excel_data.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Convert DataFrame to JSON-serializable format
            sheet_data = {
                "columns": df.columns.tolist(),
                "data": df.to_dict('records'),
                "shape": df.shape,
                "dtypes": df.dtypes.astype(str).to_dict()
            }
            
            # Add basic statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                sheet_data["statistics"] = df[numeric_cols].describe().to_dict()
            
            sheets_data[sheet_name] = sheet_data
        
        return {
            "content_type": "excel",
            "sheets": sheets_data,
            "sheet_names": list(sheets_data.keys()),
            "total_sheets": len(sheets_data)
        }
    
    def _process_csv(self, file_path: str) -> Dict[str, Any]:
        """Process CSV files"""
        data = []
        
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'iso-8859-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding, newline='') as csvfile:
                    # Detect delimiter
                    sample = csvfile.read(1024)
                    csvfile.seek(0)
                    sniffer = csv.Sniffer()
                    delimiter = sniffer.sniff(sample).delimiter
                    
                    reader = csv.DictReader(csvfile, delimiter=delimiter)
                    data = list(reader)
                    columns = list(reader.fieldnames) if reader.fieldnames else []
                    break
            except UnicodeDecodeError:
                continue
        
        if not data:
            raise Exception("Could not decode CSV file with any common encoding")
        
        # Convert to pandas for analysis if available
        if PANDAS_AVAILABLE:
            df = pd.DataFrame(data)
            numeric_cols = df.select_dtypes(include=['number']).columns
            statistics = df[numeric_cols].describe().to_dict() if len(numeric_cols) > 0 else {}
        else:
            statistics = {}
        
        return {
            "content_type": "csv",
            "columns": columns,
            "data": data,
            "total_rows": len(data),
            "total_columns": len(columns),
            "statistics": statistics
        }
    
    def _process_text(self, file_path: str) -> Dict[str, Any]:
        """Process plain text files"""
        encodings = ['utf-8', 'utf-8-sig', 'iso-8859-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                    break
            except UnicodeDecodeError:
                continue
        else:
            raise Exception("Could not decode text file with any common encoding")
        
        lines = content.split('\n')
        
        return {
            "content_type": "text",
            "content": content,
            "lines": lines,
            "total_lines": len(lines),
            "total_characters": len(content),
            "total_words": len(content.split())
        }
    
    def _process_json(self, file_path: str) -> Dict[str, Any]:
        """Process JSON files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            json_data = json.load(file)
        
        return {
            "content_type": "json",
            "data": json_data,
            "data_type": type(json_data).__name__,
            "keys": list(json_data.keys()) if isinstance(json_data, dict) else None,
            "length": len(json_data) if isinstance(json_data, (list, dict)) else None
        }
    
    def _process_xml(self, file_path: str) -> Dict[str, Any]:
        """Process XML files"""
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        def xml_to_dict(element):
            result = {}
            
            # Add attributes
            if element.attrib:
                result['@attributes'] = element.attrib
            
            # Add text content
            if element.text and element.text.strip():
                result['text'] = element.text.strip()
            
            # Add children
            children = {}
            for child in element:
                child_data = xml_to_dict(child)
                if child.tag in children:
                    if not isinstance(children[child.tag], list):
                        children[child.tag] = [children[child.tag]]
                    children[child.tag].append(child_data)
                else:
                    children[child.tag] = child_data
            
            result.update(children)
            return result
        
        xml_dict = xml_to_dict(root)
        
        return {
            "content_type": "xml",
            "root_tag": root.tag,
            "data": {root.tag: xml_dict},
            "attributes": root.attrib
        }
    
    def _process_html(self, file_path: str) -> Dict[str, Any]:
        """Process HTML files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Extract text content
        text_content = soup.get_text()
        
        # Extract structured data
        headings = []
        for i in range(1, 7):
            headings.extend([h.get_text().strip() for h in soup.find_all(f'h{i}')])
        
        links = [{'text': a.get_text().strip(), 'href': a.get('href')} 
                for a in soup.find_all('a', href=True)]
        
        images = [{'alt': img.get('alt', ''), 'src': img.get('src')} 
                 for img in soup.find_all('img')]
        
        return {
            "content_type": "html",
            "title": soup.title.string if soup.title else None,
            "text_content": text_content,
            "headings": headings,
            "links": links,
            "images": images,
            "meta_tags": [{'name': meta.get('name'), 'content': meta.get('content')} 
                         for meta in soup.find_all('meta') if meta.get('name')]
        }
    
    def _process_markdown(self, file_path: str) -> Dict[str, Any]:
        """Process Markdown files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Convert to HTML for structured parsing
        html_content = markdown.markdown(content)
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Extract headings
        headings = []
        for i in range(1, 7):
            headings.extend([h.get_text().strip() for h in soup.find_all(f'h{i}')])
        
        return {
            "content_type": "markdown",
            "raw_content": content,
            "html_content": html_content,
            "text_content": soup.get_text(),
            "headings": headings
        }
    
    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process PowerPoint presentations"""
        prs = Presentation(file_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": i + 1,
                "title": "",
                "content": [],
                "notes": ""
            }
            
            # Extract shapes and text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.shape_type == 1:  # Title
                        slide_content["title"] = shape.text
                    else:
                        slide_content["content"].append(shape.text)
            
            # Extract notes
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_content["notes"] = slide.notes_slide.notes_text_frame.text
            
            slides_data.append(slide_content)
        
        return {
            "content_type": "pptx",
            "slides": slides_data,
            "total_slides": len(slides_data),
            "slide_dimensions": (prs.slide_width, prs.slide_height)
        }
    
    def _process_image_ocr(self, file_path: str) -> Dict[str, Any]:
        """Process images with OCR"""
        image = Image.open(file_path)
        metadata = {
            "format": image.format,
            "mode": image.mode,
            "size": image.size,
            "has_transparency": "transparency" in image.info
        }

        # 1) Try Tesseract if available
        if OCR_AVAILABLE:
            try:
                extracted_text = pytesseract.image_to_string(image)
                return {
                    "content_type": "image_ocr",
                    "extracted_text": extracted_text.strip(),
                    "image_metadata": metadata,
                    "engine": "tesseract"
                }
            except Exception:
                pass

        # 2) Fallback to transformer-based OCR if available
        if TRANSFORMER_OCR_AVAILABLE:
            try:
                processor = TrOCRProcessor.from_pretrained("microsoft/trocr-base-printed")
                model = VisionEncoderDecoderModel.from_pretrained("microsoft/trocr-base-printed")
                pixel_values = processor(images=image, return_tensors="pt").pixel_values
                generated_ids = model.generate(pixel_values)
                extracted_text = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
                return {
                    "content_type": "image_ocr",
                    "extracted_text": extracted_text.strip(),
                    "image_metadata": metadata,
                    "engine": "trocr-base"
                }
            except Exception as e:
                raise Exception(f"Transformer OCR failed: {str(e)}")

        # 3) Try hosted Hugging Face Inference API if available/configured
        if HF_INFERENCE_AVAILABLE:
            try:
                with open(file_path, 'rb') as f:
                    image_bytes = f.read()
                hf_client = HFInferenceClient()
                hf_result = hf_client.ocr_image(image_bytes)
                if "error" not in hf_result:
                    return {
                        "content_type": "image_ocr",
                        "extracted_text": hf_result.get("text", "").strip(),
                        "image_metadata": metadata,
                        "engine": "hf_inference",
                        "hf_raw": hf_result.get("raw")
                    }
            except Exception:
                pass

        # 4) If no engine succeeded
        raise Exception(
            "OCR processing failed: No OCR engine succeeded (tesseract, local transformer, or HF Inference API)."
        )
